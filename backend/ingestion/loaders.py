from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
import requests
import pytesseract
from PIL import Image
import json
import base64
from io import BytesIO
from groq import Groq
import numpy as np
import xml.etree.ElementTree as ET

from backend.config.settings import GROQ_API_KEY, IMAGE_OCR_MODEL, TESSERACT_CMD


if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD


def _clean_metadata_value(value):
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value).strip()


def _metadata_section(title, entries):
    lines = []
    for label, value in entries:
        cleaned = _clean_metadata_value(value)
        if cleaned:
            lines.append(f"{label}: {cleaned}")

    if not lines:
        return ""

    return f"{title}\n" + "\n".join(lines)


def _office_metadata(core_properties, document_type):
    return _metadata_section(
        "Document metadata",
        [
            ("Document type", document_type),
            ("Title", getattr(core_properties, "title", None)),
            ("Subject", getattr(core_properties, "subject", None)),
            ("Author / creator / made by", getattr(core_properties, "author", None)),
            ("Last modified by", getattr(core_properties, "last_modified_by", None)),
            ("Created", getattr(core_properties, "created", None)),
            ("Modified", getattr(core_properties, "modified", None)),
            ("Keywords", getattr(core_properties, "keywords", None)),
            ("Category", getattr(core_properties, "category", None)),
            ("Comments", getattr(core_properties, "comments", None)),
        ],
    )


def _rapidocr_image(image):
    try:
        from rapidocr_onnxruntime import RapidOCR
    except Exception as e:
        raise RuntimeError("rapidocr-onnxruntime is not installed.") from e

    engine = RapidOCR()
    image_array = np.array(image.convert("RGB"))
    result, _ = engine(image_array)
    if not result:
        return ""

    lines = []
    for item in result:
        if len(item) >= 2:
            lines.append(str(item[1]).strip())

    return "\n".join(line for line in lines if line)


def _groq_vision_ocr(image):
    if not GROQ_API_KEY:
        raise RuntimeError("GROQ_API_KEY is missing. Cannot run vision OCR fallback.")

    buffer = BytesIO()
    image.save(buffer, format="PNG")
    image_b64 = base64.b64encode(buffer.getvalue()).decode("ascii")

    client = Groq(api_key=GROQ_API_KEY)
    response = client.chat.completions.create(
        model=IMAGE_OCR_MODEL,
        temperature=0,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all readable text from this image. Return plain text only."},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}"}},
                ],
            }
        ],
    )
    return (response.choices[0].message.content or "").strip()


def _ocr_image(image):
    errors = []

    try:
        text = pytesseract.image_to_string(image)
        if text and text.strip():
            return text
        errors.append("Tesseract returned empty text.")
    except Exception as e:
        errors.append(f"Tesseract failed: {e}")

    try:
        text = _rapidocr_image(image)
        if text and text.strip():
            return text
        errors.append("RapidOCR returned empty text.")
    except Exception as e:
        errors.append(f"RapidOCR failed: {e}")

    try:
        text = _groq_vision_ocr(image)
        if text and text.strip():
            return text
        errors.append("Vision OCR returned empty text.")
    except Exception as e:
        errors.append(f"Vision OCR failed: {e}")

    raise RuntimeError(
        "OCR could not extract text. "
        "Install/configure Tesseract, keep rapidocr-onnxruntime installed, or set GROQ_API_KEY for vision OCR. "
        + " | ".join(errors)
    )


def _load_pdf_with_pypdf(file):
    from pypdf import PdfReader

    file.seek(0)
    reader = PdfReader(file)
    text = []
    metadata = getattr(reader, "metadata", None)
    if metadata:
        text.append(
            _metadata_section(
                "Document metadata",
                [
                    ("Document type", "PDF document"),
                    ("Title", metadata.get("/Title")),
                    ("Author / creator / made by", metadata.get("/Author")),
                    ("Subject", metadata.get("/Subject")),
                    ("Creator application", metadata.get("/Creator")),
                    ("Producer", metadata.get("/Producer")),
                    ("Created", metadata.get("/CreationDate")),
                    ("Modified", metadata.get("/ModDate")),
                    ("Keywords", metadata.get("/Keywords")),
                ],
            )
        )
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            text.append(page_text)
    return "\n".join(text)


def _open_pdf(file):
    try:
        import fitz
    except ImportError as e:
        raise RuntimeError(
            "PDF extraction is unavailable because PyMuPDF could not be loaded. "
            "On this machine, Windows Application Control appears to be blocking "
            "PyMuPDF's native DLL. Allow or reinstall PyMuPDF, or use another file type."
        ) from e

    return fitz.open(stream=file.read(), filetype="pdf")


def _load_pdf_with_pymupdf(file):
    file.seek(0)
    with _open_pdf(file) as doc:
        return "\n".join(page.get_text() for page in doc if page.get_text().strip())


def _load_pdf_with_ocr(file):
    import fitz

    file.seek(0)
    text = []
    errors = []
    with _open_pdf(file) as doc:
        for page_number, page in enumerate(doc, start=1):
            try:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image = Image.open(BytesIO(pixmap.tobytes("png")))
                page_text = _ocr_image(image)
                if page_text.strip():
                    text.append(f"Page {page_number}:\n{page_text.strip()}")
            except Exception as e:
                errors.append(f"page {page_number}: {e}")

    if text:
        return "\n\n".join(text)

    if errors:
        raise RuntimeError("PDF OCR failed for all pages. " + " | ".join(errors))

    return ""


def load_pdf(file):
    errors = []

    try:
        text = _load_pdf_with_pypdf(file)
        if text.strip():
            return text
        errors.append("pypdf returned empty text.")
    except Exception as e:
        errors.append(f"pypdf failed: {e}")

    try:
        text = _load_pdf_with_pymupdf(file)
        if text.strip():
            return text
        errors.append("PyMuPDF returned empty text.")
    except Exception as e:
        errors.append(f"PyMuPDF text extraction failed: {e}")

    try:
        text = _load_pdf_with_ocr(file)
        if text.strip():
            return text
        errors.append("PDF OCR returned empty text.")
    except Exception as e:
        errors.append(f"PDF OCR failed: {e}")

    raise RuntimeError(
        "PDF text extraction returned no text. The PDF may be scanned, image-only, or unreadable. "
        + " | ".join(errors)
    )

def load_docx(file):
    file.seek(0)
    doc = Document(file)
    parts = [_office_metadata(doc.core_properties, "Word document")]
    parts.extend(p.text for p in doc.paragraphs if p.text and p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(part for part in parts if part and part.strip())


def load_pptx(file):
    file.seek(0)
    prs = Presentation(file)
    text = [_office_metadata(prs.core_properties, "PowerPoint presentation")]
    for slide_number, s in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in s.shapes:
            slide_text.extend(_extract_pptx_shape_text(shape))
        if getattr(s, "has_notes_slide", False):
            notes = getattr(s.notes_slide, "notes_text_frame", None)
            notes_text = str(getattr(notes, "text", "") or "").strip()
            if notes_text:
                slide_text.append(f"Speaker notes:\n{notes_text}")
        if slide_text:
            text.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))
    return "\n\n".join(part for part in text if part and part.strip())


def _extract_pptx_shape_text(shape):
    parts = []

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    if hasattr(shape, "text"):
        shape_text = str(shape.text or "").strip()
        if shape_text:
            parts.append(shape_text)

    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            parts.extend(_extract_pptx_shape_text(child))

    return parts


def load_csv(file):
    file.seek(0)
    return pd.read_csv(file).to_string()


def load_image(file):
    file.seek(0)
    image = Image.open(file)
    return _ocr_image(image)

def load_json(file):
    file.seek(0)
    return json.dumps(json.load(file), indent=2, ensure_ascii=False)


def load_text(file):
    file.seek(0)
    content = file.read()
    if isinstance(content, bytes):
        return content.decode("utf-8", errors="replace")
    return content


def load_html(file):
    soup = BeautifulSoup(load_text(file), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator=" ")


def load_excel(file):
    file.seek(0)
    sheets = pd.read_excel(file, sheet_name=None)
    text = []
    for sheet_name, dataframe in sheets.items():
        text.append(f"Sheet: {sheet_name}")
        text.append(dataframe.to_string(index=False))
    return "\n\n".join(text)


def load_xml(file):
    content = load_text(file)
    root = ET.fromstring(content)
    texts = []
    for element in root.iter():
        if element.text and element.text.strip():
            texts.append(element.text.strip())
    return "\n".join(texts)


def load_url(url, return_response=False):
    try:
        res = requests.get(url, timeout=10)

        if res.status_code != 200:
            return None if return_response else ""

        if return_response:
            return res

        soup = BeautifulSoup(res.text, "html.parser")

        # remove scripts/styles
        for tag in soup(["script", "style"]):
            tag.decompose()

        text = soup.get_text(separator=" ")

        return text

    except Exception:
        return None if return_response else ""
